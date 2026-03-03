#!/usr/bin/env python3
"""
Script to generate PowerPoint presentation for IAM Gateway UPEC project.
Run via Docker: docker run --rm -v $(pwd)/presentations:/presentations python:3.11-slim \
    bash -c "pip install python-pptx -q && python3 /presentations/generate_pptx.py"
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os


def rgb(r, g, b):
    """Create an RGBColor from r,g,b values."""
    from pptx.dml.color import RGBColor
    return RGBColor(r, g, b)


# ── Color palette ──
DARK_BG      = rgb(15, 23, 42)
DARK_CARD    = rgb(30, 41, 59)
PRIMARY      = rgb(59, 130, 246)
PRIMARY_DARK = rgb(30, 64, 175)
ACCENT       = rgb(6, 182, 212)
SUCCESS      = rgb(16, 185, 129)
WARNING      = rgb(245, 158, 11)
DANGER       = rgb(239, 68, 68)
PURPLE       = rgb(124, 58, 237)
WHITE        = rgb(255, 255, 255)
LIGHT        = rgb(203, 213, 225)
GRAY         = rgb(148, 163, 184)
MUTED        = rgb(100, 116, 139)

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
LOGO_IUT   = os.path.join(SCRIPT_DIR, "cropped-Logo-INFO-FOND-BLANC.jpg")
LOGO_LISSI = os.path.join(SCRIPT_DIR, "Lissi-cmjn.png")
SCREENSHOTS = os.path.join(SCRIPT_DIR, "screenshots")


# ════════════════════════════════════════════════════════════════
# Helper functions
# ════════════════════════════════════════════════════════════════

def _add_bg(slide, prs):
    """Dark background + top accent line."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = PRIMARY
    line.line.fill.background()


def _add_slide_title(slide, text, subtitle=None):
    """Add title (and optional subtitle) to a slide."""
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.25), Inches(8.8), Inches(0.65))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.6), Inches(0.85), Inches(8.8), Inches(0.4))
        tf2 = sub.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.color.rgb = ACCENT


def _add_bullet_box(slide, left, top, width, height, items):
    """Add a text box with bullet items. Supports ## headers and - bullets."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        if item == "":
            p.text = ""
            p.space_after = Pt(4)
        elif item.startswith("## "):
            p.text = item[3:]
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = ACCENT
            p.space_before = Pt(10)
        elif item.startswith("- "):
            p.text = "  " + item
            p.font.size = Pt(14)
            p.font.color.rgb = LIGHT
        elif item.startswith(">>> "):
            p.text = item[4:]
            p.font.size = Pt(14)
            p.font.color.rgb = WARNING
            p.font.italic = True
        else:
            p.text = item
            p.font.size = Pt(15)
            p.font.color.rgb = WHITE


def add_content_slide(prs, title, left_items, right_items=None, image_path=None, subtitle=None):
    """Generic content slide with left content + optional right content/image."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, prs)
    _add_slide_title(slide, title, subtitle)

    top = Inches(1.3)
    h = Inches(5.5)

    if right_items or image_path:
        _add_bullet_box(slide, Inches(0.5), top, Inches(4.5), h, left_items)
        if image_path and os.path.exists(image_path):
            slide.shapes.add_picture(image_path, Inches(5.2), top, width=Inches(4.5))
        elif right_items:
            _add_bullet_box(slide, Inches(5.2), top, Inches(4.3), h, right_items)
    else:
        _add_bullet_box(slide, Inches(0.5), top, Inches(9), h, left_items)

    return slide


def add_table_slide(prs, title, headers, rows, subtitle=None):
    """Slide with a styled table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, prs)
    _add_slide_title(slide, title, subtitle)

    cols = len(headers)
    row_count = len(rows) + 1
    row_height = min(Inches(0.45), Inches(5.0) / row_count)
    table_height = row_height * row_count

    tbl = slide.shapes.add_table(
        row_count, cols, Inches(0.5), Inches(1.3), Inches(9), table_height
    ).table

    # Header row
    for i, h in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY_DARK
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE

    # Data rows
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = DARK_CARD if ri % 2 == 0 else DARK_BG
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = WHITE

    return slide


def add_two_image_slide(prs, title, img_left, img_right, caption_left="", caption_right="", subtitle=None):
    """Slide with two images side by side."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, prs)
    _add_slide_title(slide, title, subtitle)

    top = Inches(1.4)
    img_w = Inches(4.3)

    if os.path.exists(img_left):
        slide.shapes.add_picture(img_left, Inches(0.4), top, width=img_w)
    if os.path.exists(img_right):
        slide.shapes.add_picture(img_right, Inches(5.2), top, width=img_w)

    if caption_left:
        cb = slide.shapes.add_textbox(Inches(0.4), Inches(6.0), img_w, Inches(0.4))
        p = cb.text_frame.paragraphs[0]
        p.text = caption_left
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER

    if caption_right:
        cb2 = slide.shapes.add_textbox(Inches(5.2), Inches(6.0), img_w, Inches(0.4))
        p2 = cb2.text_frame.paragraphs[0]
        p2.text = caption_right
        p2.font.size = Pt(12)
        p2.font.color.rgb = GRAY
        p2.alignment = PP_ALIGN.CENTER

    return slide


def add_progress_table(prs, title, items, subtitle=None):
    """Table with feature name, %, description. Colored % column."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, prs)
    _add_slide_title(slide, title, subtitle)

    headers = ["Fonctionnalite", "%", "Details"]
    cols = 3
    row_count = len(items) + 1

    tbl = slide.shapes.add_table(
        row_count, cols, Inches(0.5), Inches(1.3), Inches(9), Inches(0.45 * row_count)
    ).table

    # Set column widths
    tbl.columns[0].width = Inches(3.5)
    tbl.columns[1].width = Inches(1.0)
    tbl.columns[2].width = Inches(4.5)

    for i, h in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY_DARK
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE

    for ri, (name, pct, detail) in enumerate(items):
        # Name
        c0 = tbl.cell(ri + 1, 0)
        c0.text = name
        c0.fill.solid()
        c0.fill.fore_color.rgb = DARK_CARD if ri % 2 == 0 else DARK_BG
        c0.text_frame.paragraphs[0].font.size = Pt(12)
        c0.text_frame.paragraphs[0].font.color.rgb = WHITE
        c0.text_frame.paragraphs[0].font.bold = True

        # Percentage
        c1 = tbl.cell(ri + 1, 1)
        c1.text = pct
        c1.fill.solid()
        pct_val = int(pct.replace("%", "").strip())
        if pct_val == 100:
            c1.fill.fore_color.rgb = rgb(5, 46, 22)
        elif pct_val >= 70:
            c1.fill.fore_color.rgb = rgb(30, 58, 12)
        elif pct_val >= 40:
            c1.fill.fore_color.rgb = rgb(66, 32, 6)
        else:
            c1.fill.fore_color.rgb = rgb(69, 10, 10)
        p1 = c1.text_frame.paragraphs[0]
        p1.font.size = Pt(13)
        p1.font.bold = True
        p1.alignment = PP_ALIGN.CENTER
        if pct_val == 100:
            p1.font.color.rgb = SUCCESS
        elif pct_val >= 70:
            p1.font.color.rgb = rgb(134, 239, 172)
        elif pct_val >= 40:
            p1.font.color.rgb = WARNING
        else:
            p1.font.color.rgb = DANGER

        # Detail
        c2 = tbl.cell(ri + 1, 2)
        c2.text = detail
        c2.fill.solid()
        c2.fill.fore_color.rgb = DARK_CARD if ri % 2 == 0 else DARK_BG
        c2.text_frame.paragraphs[0].font.size = Pt(11)
        c2.text_frame.paragraphs[0].font.color.rgb = LIGHT

    return slide


# ════════════════════════════════════════════════════════════════
# Main presentation
# ════════════════════════════════════════════════════════════════

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ─────────────────────────────────────────────
    # SLIDE 1 - Title
    # ─────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()

    # Logos
    if os.path.exists(LOGO_IUT):
        slide.shapes.add_picture(LOGO_IUT, Inches(0.5), Inches(0.3), height=Inches(1.2))
    if os.path.exists(LOGO_LISSI):
        slide.shapes.add_picture(LOGO_LISSI, Inches(6.5), Inches(0.3), height=Inches(1.0))

    # Accent line
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(2.2), Inches(6), Inches(0.04))
    ln.fill.solid()
    ln.fill.fore_color.rgb = PRIMARY
    ln.line.fill.background()

    # Title
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = "IAM Gateway"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sb = slide.shapes.add_textbox(Inches(0.5), Inches(3.4), Inches(9), Inches(0.6))
    p2 = sb.text_frame.paragraphs[0]
    p2.text = "Passerelle de Provisionnement Intelligente des Identites"
    p2.font.size = Pt(22)
    p2.font.color.rgb = ACCENT
    p2.alignment = PP_ALIGN.CENTER

    # Authors
    ab = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(0.5))
    p3 = ab.text_frame.paragraphs[0]
    p3.text = "Equipe MOE : Zhmuryk Andrii  &  Aydin Ibrahim"
    p3.font.size = Pt(20)
    p3.font.bold = True
    p3.font.color.rgb = WHITE
    p3.alignment = PP_ALIGN.CENTER

    # University info
    ub = slide.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(9), Inches(1.5))
    tf = ub.text_frame
    lines = [
        "Product Owner & Architecture : A. CHIBANI",
        "BUT Informatique 3e annee - UPEC",
        "IUT Creteil-Vitry  |  Laboratoire LISSI",
        "NEXUS AI Innovation Lab  -  Fevrier 2026",
    ]
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(15)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER

    # ─────────────────────────────────────────────
    # SLIDE 2 - Sommaire
    # ─────────────────────────────────────────────
    add_content_slide(prs, "Sommaire", [
        "## 1.  Problematique & Contexte",
        "## 2.  Notre Solution",
        "## 3.  Architecture & Stack Technique",
        "## 4.  Backlog des Fonctionnalites",
        "## 5.  Fonctionnalites Implementees + Screenshots",
        "## 6.  Fonctionnalites en Production",
        "## 7.  Fonctionnalites a Finaliser (% progression)",
        "## 8.  Fonctionnalites Futures",
        "## 9.  Environnement de Developpement",
        "## 10. Difficultes Rencontrees",
        "## 11. Conclusion & Competences Acquises",
    ])

    # ─────────────────────────────────────────────
    # SLIDE 3 - Problematique
    # ─────────────────────────────────────────────
    add_content_slide(prs, "Problematique & Contexte", [
        "## Le defi de la gestion des identites",
        "",
        "- 1 employe = 5 a 10 comptes (AD, LDAP, ERP, Cloud...)",
        "- Creation manuelle = erreurs + delais",
        "- Depart employe = comptes orphelins oublies",
        "- Aucune tracabilite des acces",
        "- Risques de securite majeurs",
        "",
        "## Chiffres cles",
        "- 30% des violations liees aux comptes orphelins",
        "- 4h en moyenne pour creer les acces d'un employe",
        "- 0% de tracabilite sans outil IGA",
    ], [
        "## Systemes concernes",
        "",
        "- Active Directory / LDAP",
        "- ERP (Odoo, SAP)",
        "- Applications metier (bases SQL)",
        "- Services Cloud (Keycloak, OAuth)",
        "- Intranet & outils internes",
        "",
        "## Besoin",
        ">>> Automatiser et centraliser le cycle",
        ">>> de vie des identites numeriques",
    ])

    # ─────────────────────────────────────────────
    # SLIDE 4 - Solution
    # ─────────────────────────────────────────────
    add_content_slide(prs, "Notre Solution : IAM Gateway", [
        "## Passerelle No-Code pour MidPoint",
        "",
        "- MidPoint (Evolveum) = puissant mais complexe",
        "- Notre gateway simplifie l'interface",
        "- Interface web intuitive et moderne",
        "",
        "## 3 piliers de la solution",
        "",
        "## 1. Automatisation",
        "- Provisionnement automatique multi-systemes",
        "- Regles dynamiques (Jinja2)",
        "- Workflows d'approbation multi-niveaux",
    ], [
        "",
        "## 2. Centralisation",
        "- Un seul point d'entree pour tous les systemes",
        "- Vue temps reel des identites",
        "- Reconciliation automatique des ecarts",
        "",
        "## 3. Conformite",
        "- Tracabilite complete (audit logs)",
        "- Recherche semantique IA dans les logs",
        "- Assistant IA integre",
        "",
        ">>> Innovation : rendre l'IAM accessible a tous",
    ])

    # ─────────────────────────────────────────────
    # SLIDE 5 - Architecture Schema
    # ─────────────────────────────────────────────
    add_content_slide(prs, "Schema d'Architecture Globale", [
        "## Flux de provisionnement",
        "",
        "SOURCES (Odoo RH / CSV / Utilisateur)",
        "              |",
        "              v",
        "GATEWAY IAM (FastAPI - port 8000)",
        "  [Regles] [Workflows] [IA]",
        "              |",
        "              v",
        ">>> MIDPOINT 4.4 - COEUR IGA <<<",
        "  Orchestrateur du provisionnement",
        "  Roles, Policies, Reconciliation",
        "              |",
        "              v",
        "SYSTEMES CIBLES provisionnes par MidPoint",
        "  LDAP | Keycloak | Odoo | SQL",
    ], [
        "## Pourquoi MidPoint au centre ?",
        "",
        "- Moteur IGA open source de reference",
        "- Gere les projections (shadows)",
        "- Reconciliation native entre systemes",
        "- Roles et policies centralises",
        "",
        "## Notre Gateway = No-Code pour MidPoint",
        "",
        "- Gateway = simplification (interface)",
        "- MidPoint = orchestration (provisionnement)",
        "",
        "## Infrastructure",
        "- 14 services Docker Compose",
        "- 5 bases PostgreSQL",
        "- Redis (cache) + Qdrant (vecteurs IA)",
    ])

    # ─────────────────────────────────────────────
    # SLIDE 6 - Architecture Technique
    # ─────────────────────────────────────────────
    add_table_slide(prs, "Architecture Technique - 14 Services Docker",
        ["Service", "Technologie", "Port", "Role"],
        [
            ["gateway-frontend", "React 18 / TypeScript", "3000", "Interface utilisateur"],
            ["gateway", "FastAPI / Python", "8000", "API & logique metier"],
            ["midpoint", "MidPoint 4.4", "8080", "IGA central"],
            ["openldap", "OpenLDAP 2.6", "10389", "Annuaire LDAP"],
            ["odoo", "Odoo 17", "8069", "ERP / SIRH"],
            ["keycloak", "Keycloak 23", "8081", "SSO / OIDC"],
            ["redis", "Redis 7", "6379", "Cache & sessions"],
            ["qdrant", "Qdrant 1.7", "6333", "Recherche vectorielle IA"],
            ["5x PostgreSQL", "PostgreSQL 15/16", "-", "Bases de donnees"],
            ["phpldapadmin", "phpLDAPadmin", "8443", "Admin LDAP"],
        ]
    )

    # ─────────────────────────────────────────────
    # SLIDE 7 - Stack
    # ─────────────────────────────────────────────
    add_content_slide(prs, "Stack Technologique", [
        "## Backend (Python)",
        "- FastAPI 0.109 (async)",
        "- SQLModel + PostgreSQL",
        "- JWT + OAuth2 (securite)",
        "- ldap3 (LDAP natif)",
        "- Jinja2 Sandbox (moteur regles)",
        "- APScheduler (planification)",
        "- OpenAI API (assistant IA)",
        "- Qdrant (recherche semantique)",
    ], [
        "## Frontend (TypeScript)",
        "- React 18 + Vite",
        "- React Router v6",
        "- Zustand (state management)",
        "- Tailwind CSS (styling)",
        "- i18next (FR, EN, UK)",
        "- Lucide React (icons)",
        "- Axios (HTTP client)",
        "- Recharts (graphiques)",
    ])

    # ─────────────────────────────────────────────
    # SLIDE 8 - Backlog Epic 1
    # ─────────────────────────────────────────────
    add_table_slide(prs, "Backlog - Gestion des Utilisateurs",
        ["User Story", "Priorite", "Pts", "Status"],
        [
            ["Creer un utilisateur dans tous les systemes cibles", "Must Have", "8", "Done"],
            ["Modifier les attributs d'un utilisateur", "Must Have", "5", "Done"],
            ["Desactiver/supprimer un compte partout", "Must Have", "5", "Done"],
            ["Voir l'historique des modifications", "Should Have", "3", "Done"],
            ["Assigner des roles MidPoint a un utilisateur", "Must Have", "5", "Done"],
            ["Gerer les groupes LDAP (membres)", "Must Have", "5", "Done"],
            ["Importer utilisateurs depuis CSV", "Should Have", "5", "Done"],
        ],
        subtitle="Epic 1 : Gestion du cycle de vie des identites"
    )

    # ─────────────────────────────────────────────
    # SLIDE 9 - Backlog Epic 2, 3
    # ─────────────────────────────────────────────
    add_table_slide(prs, "Backlog - Connecteurs & Regles",
        ["User Story", "Priorite", "Pts", "Status"],
        [
            ["Connecteur LDAP (OpenLDAP, AD)", "Must Have", "8", "Done"],
            ["Connecteur Odoo (ERP via XML-RPC)", "Must Have", "8", "Done"],
            ["Connecteur SQL (PostgreSQL)", "Must Have", "5", "Done"],
            ["Connecteur MidPoint (REST API)", "Must Have", "13", "Done"],
            ["Connecteur Keycloak (OIDC)", "Should Have", "8", "70%"],
            ["Connecteur CSV (import/export)", "Should Have", "5", "Done"],
            ["Moteur de regles Jinja2", "Must Have", "8", "Done"],
            ["Workflows d'approbation multi-niveaux", "Must Have", "8", "Done"],
        ],
        subtitle="Epic 2 : Connecteurs  |  Epic 3 : Regles & Workflows"
    )

    # ─────────────────────────────────────────────
    # SLIDE 10 - Backlog Epic 4, 5
    # ─────────────────────────────────────────────
    add_table_slide(prs, "Backlog - Audit, IA & Monitoring",
        ["User Story", "Priorite", "Pts", "Status"],
        [
            ["Dashboard temps reel avec metriques", "Must Have", "5", "Done"],
            ["Audit logs avec recherche full-text", "Must Have", "5", "Done"],
            ["Recherche semantique IA (Qdrant)", "Could Have", "8", "80%"],
            ["Assistant IA (chat GPT-4)", "Could Have", "8", "Done"],
            ["Reconciliation multi-systemes", "Must Have", "8", "Done"],
            ["Comparaison live entre systemes", "Should Have", "5", "Done"],
            ["Planificateur de synchronisations", "Should Have", "5", "Done"],
            ["Notifications email", "Should Have", "5", "50%"],
        ],
        subtitle="Epic 4 : Audit & Conformite  |  Epic 5 : Monitoring & IA"
    )

    # ─────────────────────────────────────────────
    # SLIDE 11 - Features implementees: Dashboard
    # ─────────────────────────────────────────────
    dash_img = os.path.join(SCREENSHOTS, "dashboard.png")
    add_content_slide(prs, "Implementee : Dashboard", [
        "## Dashboard Principal",
        "",
        "- Vue d'ensemble en temps reel",
        "- Nombre d'utilisateurs par systeme",
        "- Etat de sante des connecteurs",
        "- Dernieres operations",
        "- Services MidPoint en temps reel",
        "",
        "## API",
        "- GET /api/v1/admin/status",
        "- GET /api/v1/admin/metrics",
        "- GET /api/v1/connectors/health",
    ], image_path=dash_img,
    subtitle="100% - Operationnel en production")

    # ─────────────────────────────────────────────
    # SLIDE 12 - Features implementees: Provisioning + Connecteurs
    # ─────────────────────────────────────────────
    add_content_slide(prs, "Implementee : Provisionnement Multi-Systemes", [
        "## Operations de provisionnement",
        "",
        "- Creation utilisateur dans LDAP + Odoo + MidPoint + SQL",
        "- Modification des attributs en cascade",
        "- Desactivation/suppression multi-systemes",
        "- Selection visuelle des systemes cibles",
        "- Formulaire dynamique avec niveaux de permission",
        "",
        "## 6 connecteurs natifs",
        "- LDAP (OpenLDAP / Active Directory)",
        "- Odoo 17 (XML-RPC)",
        "- PostgreSQL (SQL direct)",
        "- MidPoint (REST API)",
        "- Keycloak (OIDC)",
        "- CSV (import/export)",
    ], [
        "## Flux d'une operation",
        "",
        "1. Saisie formulaire (attributs)",
        "2. Selection systemes cibles",
        "3. Evaluation des regles",
        "4. Workflow d'approbation",
        "5. Execution multi-systeme",
        "6. Log audit + notification",
        "",
        "## Connecteur Wizard",
        "- Formulaire 2 etapes (schema-driven)",
        "- Test de connexion en direct",
        "- Configuration JSON Schema",
    ],
    subtitle="100% - Provisionnement + 6 connecteurs operationnels")

    # ─────────────────────────────────────────────
    # SLIDE 13 - Features: MidPoint Users + Roles
    # ─────────────────────────────────────────────
    mp_users = os.path.join(SCREENSHOTS, "midpoint-users.png")
    mp_roles = os.path.join(SCREENSHOTS, "midpoint-roles.png")
    add_two_image_slide(prs,
        "Implementee : Gestion MidPoint",
        mp_users, mp_roles,
        "Utilisateurs MidPoint", "Roles MidPoint",
        subtitle="100% - CRUD complet + attribution roles"
    )

    # ─────────────────────────────────────────────
    # SLIDE 14 - Features: Comparaison Live + Reconciliation
    # ─────────────────────────────────────────────
    comp_img = os.path.join(SCREENSHOTS, "comparaison-live.png")
    recon_img = os.path.join(SCREENSHOTS, "reconciliation.png")
    add_two_image_slide(prs,
        "Implementee : Comparaison & Reconciliation",
        comp_img, recon_img,
        "Comparaison live entre systemes", "Reconciliation des divergences",
        subtitle="100% - Detection et resolution des ecarts"
    )

    # ─────────────────────────────────────────────
    # SLIDE 15 - Features: LDAP Groups
    # ─────────────────────────────────────────────
    ldap_img = os.path.join(SCREENSHOTS, "ldap-groups.png")
    add_content_slide(prs, "Implementee : Groupes LDAP", [
        "## Gestion complete des groupes LDAP",
        "",
        "- Liste de tous les groupes",
        "- Visualisation des membres",
        "- Ajout avec autocomplete intelligent",
        "- Suppression de membres",
        "- Support groupOfUniqueNames",
        "- Filtrage automatique OIDs MidPoint",
        "",
        "## Interface",
        "- Recherche temps reel",
        "- Autocompletion depuis annuaire",
    ], image_path=ldap_img,
    subtitle="100% - Gestion groupes LDAP operationnelle")

    # ─────────────────────────────────────────────
    # SLIDE 16 - Features: Workflows + Audit + IA
    # ─────────────────────────────────────────────
    add_content_slide(prs, "Implementees : Workflows, Audit & IA", [
        "## Workflows d'Approbation",
        "- Configuration 5 niveaux",
        "- Approbation / rejet avec commentaires",
        "- Timeout configurable (72h par defaut)",
        "- Notifications par email",
        "- Historique complet des decisions",
        "",
        "## Moteur de Regles",
        "- Expressions Jinja2 sandboxees",
        "- Types: mapping, calcul, validation, aggregation",
        "- Test avec donnees fictives",
        "- Filtres integres (normalize_name, etc.)",
    ], [
        "## Audit & Conformite",
        "- Logging structure JSON",
        "- Recherche full-text dans les logs",
        "- Recherche semantique IA (Qdrant)",
        "- Filtrage par date, type, source",
        "- Conformite RGPD",
        "",
        "## Assistant IA",
        "- Chat intelligent (GPT-4)",
        "- Analyse des identites",
        "- Suggestions de regles",
        "- Aide au diagnostic",
    ],
    subtitle="100% - 4 modules operationnels")

    # ─────────────────────────────────────────────
    # SLIDE 17 - Fonctionnalites en Production
    # ─────────────────────────────────────────────
    add_table_slide(prs, "Fonctionnalites en Production",
        ["Module", "Pages UI", "Endpoints API", "Status"],
        [
            ["Dashboard", "1", "3", "Operationnel"],
            ["Provisionnement", "1 + modal", "6", "Operationnel"],
            ["Connecteurs", "1 + wizard", "5", "Operationnel"],
            ["Utilisateurs MidPoint", "1 + modal roles", "8", "Operationnel"],
            ["Utilisateurs Gateway", "1 + approbation", "5", "Operationnel"],
            ["Comparaison Live", "1", "4", "Operationnel"],
            ["Reconciliation", "1", "3", "Operationnel"],
            ["Groupes LDAP", "1", "4", "Operationnel"],
            ["Workflows", "1", "4", "Operationnel"],
            ["Audit & Logs", "1", "3", "Operationnel"],
            ["Assistant IA", "1", "2", "Operationnel"],
        ],
        subtitle="12 modules deployes et fonctionnels"
    )

    # ─────────────────────────────────────────────
    # SLIDE 18 - Fonctionnalites a finaliser (% progression)
    # ─────────────────────────────────────────────
    add_progress_table(prs, "Fonctionnalites a Finaliser",
        [
            ("Integration Keycloak (sync groupes)", "70%", "Mapping roles OIDC <-> MidPoint"),
            ("Notifications Email (templates HTML)", "50%", "Templates Jinja2 + SMTP config"),
            ("Recherche Vectorielle (refresh)", "80%", "Refresh auto des embeddings Qdrant"),
            ("Rollback Operations (nettoyage)", "60%", "Annulation en cascade + etats"),
            ("Editeur de Regles (UI)", "30%", "Formulaire creation/edition regles"),
            ("Persistance Operations (DB)", "40%", "Migrer du memory_store vers PostgreSQL"),
        ],
        subtitle="6 fonctionnalites partiellement implementees"
    )

    # ─────────────────────────────────────────────
    # SLIDE 19 - Fonctionnalites futures
    # ─────────────────────────────────────────────
    add_content_slide(prs, "Fonctionnalites a Implementer (Futur)", [
        "## Priorite Haute",
        "- Deploiement Kubernetes / Helm charts",
        "- Masquage donnees PII (RGPD)",
        "- Migration schema DB (Alembic)",
        "- Connecteur Azure AD / Entra ID",
        "",
        "## Priorite Moyenne",
        "- Connecteur Google Workspace",
        "- Connecteur Firebase",
        "- Detection Separation of Duties (SoD)",
        "- Constructeur visuel de workflows (drag & drop)",
    ], [
        "## Priorite Basse",
        "- Application mobile (React Native)",
        "- Multi-tenancy (mode SaaS)",
        "- Machine Learning (detection anomalies)",
        "- MFA biometrique",
        "",
        "## Estimation",
        "- ~60 story points restants",
        "- 3 sprints supplementaires estimes",
        "",
        ">>> Feuille de route disponible",
        ">>> sur GitHub Issues",
    ],
    subtitle="Roadmap produit post-SAE")

    # ─────────────────────────────────────────────
    # SLIDE 20 - Environnement Dev: VSCode & GitHub
    # ─────────────────────────────────────────────
    add_content_slide(prs, "Environnement de Developpement", [
        "## IDE : Visual Studio Code",
        "- Extensions : Python, ESLint, Prettier, Docker",
        "- Terminal integre pour docker-compose",
        "- Debugger FastAPI (launch.json)",
        "- Live reload frontend (Vite HMR)",
        "",
        "## Controle de version : Git & GitHub",
        "- Organisation : NEXUS-AI-Innovation-lab",
        "- Repo : IAM-Gateway",
        "- Branches : main + feature branches",
        "- Commits conventionnels",
    ], [
        "## GitHub Issues & Project",
        "- 9+ issues creees (7 closed, 2 open)",
        "- Labels : bug, enhancement, documentation",
        "- Milestones par sprint",
        "",
        "## CI/CD & Docker",
        "- docker-compose.yml (14 services)",
        "- .env pour configuration",
        "- Volumes persistants",
        "- Build frontend : npm run build",
        "",
        "## Documentation",
        "- GUIDE_DEVELOPPEUR.md (1300+ lignes)",
        "- README.md par module",
    ],
    subtitle="Outils, workflow et organisation du projet")

    # ─────────────────────────────────────────────
    # SLIDE 21 - GitHub (Issues, Commits)
    # ─────────────────────────────────────────────
    add_content_slide(prs, "GitHub : Issues, Commits & Organisation", [
        "## Organisation GitHub",
        "- NEXUS-AI-Innovation-lab/IAM-Gateway",
        "- Visibilite publique",
        "",
        "## Issues creees",
        "- #1 Configurer CI pipeline",
        "- #2 Ajouter tests unitaires backend",
        "- #3 Optimiser requetes LDAP",
        "- #4 Ameliorer UX mobile",
        "- #5 Documenter API REST (Swagger)",
        "- #6 Ajouter monitoring Prometheus",
        "- #7 Audit de securite",
    ], [
        "## Commits recents",
        "- Add i18n support (FR, EN, UK)",
        "- Add developer guide & code docs",
        "- Update poster and presentation",
        "- Add real IUT and LISSI logos",
        "- Fix user management & approval system",
        "- Add IGA connector types",
        "",
        "## Statistiques",
        "- 50+ commits",
        "- 2 contributeurs",
        "- 21,000+ lignes de code",
    ],
    subtitle="Historique et gestion du projet")

    # ─────────────────────────────────────────────
    # SLIDE 22 - Difficultes
    # ─────────────────────────────────────────────
    add_content_slide(prs, "Difficultes Rencontrees", [
        "## 1. Complexite MidPoint",
        "- API XML complexe, peu documentee",
        "- Solution : abstraction via REST API",
        "",
        "## 2. Transactions distribuees",
        "- 1 operation = 4 systemes cibles",
        "- Solution : pattern Saga + compensations",
        "",
        "## 3. Securite moteur regles",
        "- Jinja2 peut executer du code arbitraire",
        "- Solution : SandboxedEnvironment",
        "",
        "## 4. Gestion des secrets",
        "- 20+ variables sensibles",
        "- Solution : .env + validation au demarrage",
    ], [
        "## 5. Synchronisation temps reel",
        "- Coherence eventuelle entre systemes",
        "- Solution : cache invalidation + polling",
        "",
        "## 6. Mapping attributs",
        "- Schemas differents par systeme",
        "- Solution : config mappings par connecteur",
        "",
        "## 7. Couts & latence IA",
        "- OpenAI API = cout + latence",
        "- Solution : cache + rate limiting",
        "",
        "## 8. Tests d'integration",
        "- 14 services = complexite",
        "- Solution : mocks + env dedie",
    ])

    # ─────────────────────────────────────────────
    # SLIDE 23 - Conclusion & Competences
    # ─────────────────────────────────────────────
    add_content_slide(prs, "Conclusion & Competences Acquises", [
        "## Ce que nous avons realise",
        "- 12+ modules operationnels",
        "- 14 services Docker orchestres",
        "- 80+ endpoints API REST",
        "- 15 pages frontend",
        "- 3 langues (FR, EN, UK)",
        "- 16,000+ lignes Python + 5,000+ lignes TS",
        "",
        "## Competences Techniques",
        "- Architecture microservices",
        "- API REST avancee (FastAPI)",
        "- React moderne (hooks, Zustand, Router)",
        "- Docker Compose & orchestration",
        "- Securite (JWT, RBAC, OAuth2)",
        "- Bases de donnees (PostgreSQL, LDAP, Redis)",
    ], [
        "## Competences IAM / IGA",
        "- Provisionnement des identites",
        "- Reconciliation multi-systemes",
        "- Workflows d'approbation",
        "- Audit et conformite (RGPD)",
        "- RBAC / ABAC",
        "- SSO / OIDC / Keycloak",
        "- MidPoint administration",
        "",
        "## Competences Projet",
        "- Methodologie SCRUM",
        "- Gestion backlog (user stories)",
        "- Git & GitHub (issues, branches)",
        "- Documentation technique",
        "- Presentation et communication",
    ],
    subtitle="Bilan du projet SAE 2026")

    # ─────────────────────────────────────────────
    # SLIDE 24 - Merci / Questions
    # ─────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()

    # Logos
    if os.path.exists(LOGO_IUT):
        slide.shapes.add_picture(LOGO_IUT, Inches(0.5), Inches(0.3), height=Inches(1.0))
    if os.path.exists(LOGO_LISSI):
        slide.shapes.add_picture(LOGO_LISSI, Inches(6.5), Inches(0.3), height=Inches(0.9))

    # Thank you
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = "Merci de votre attention !"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    sb = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.6))
    p2 = sb.text_frame.paragraphs[0]
    p2.text = "Questions ?"
    p2.font.size = Pt(32)
    p2.font.color.rgb = ACCENT
    p2.alignment = PP_ALIGN.CENTER

    # Contact info
    cb = slide.shapes.add_textbox(Inches(1), Inches(5.0), Inches(8), Inches(2))
    tf = cb.text_frame
    contact = [
        "GitHub : github.com/NEXUS-AI-Innovation-lab/IAM-Gateway",
        "Equipe MOE : Zhmuryk Andrii & Aydin Ibrahim",
        "Product Owner : A. CHIBANI",
        "IUT Creteil-Vitry  |  Laboratoire LISSI  |  UPEC",
    ]
    for i, line in enumerate(contact):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER

    # ── Save ──
    output = os.path.join(SCRIPT_DIR, "IAM_Gateway_UPEC.pptx")
    prs.save(output)
    print(f"Presentation saved: {output}")
    print(f"Total slides: {len(prs.slides)}")
    return output


if __name__ == "__main__":
    create_presentation()
